"""
HalluCheck v3 — PowerPoint Presentation Generator
Run: python create_ppt.py
Output: HalluCheck_v3_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a  slide background
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)   # #38bdf8  sky-400 (headings)
GREEN      = RGBColor(0x34, 0xD3, 0x99)   # #34d399  verified / good
RED        = RGBColor(0xF8, 0x71, 0x71)   # #f87171  blocked / bad
YELLOW     = RGBColor(0xFB, 0xBF, 0x24)   # #fbbf24  flagged / warning
GRAY       = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8  secondary text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PANEL      = RGBColor(0x1E, 0x29, 0x3B)   # #1e293b  card/panel bg
LIGHT_TEXT = RGBColor(0xCB, 0xD5, 0xE1)   # #cbd5e1  body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # dark background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide


def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=LIGHT_TEXT,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def heading_bar(slide, title, subtitle=None):
    """Dark accent bar across the top of a content slide."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PANEL
    bar.line.color.rgb = ACCENT
    bar.line.width = Pt(2)
    txt(slide, title,
        Inches(0.4), Inches(0.1), Inches(10), Inches(0.75),
        size=32, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.4), Inches(0.8), Inches(12), Inches(0.5),
            size=14, color=GRAY)


def bullet_block(slide, items, left, top, width, item_height=Inches(0.38),
                 size=16, color=LIGHT_TEXT, marker="▸ "):
    y = top
    for item in items:
        txt(slide, f"{marker}{item}", left, y, width, item_height,
            size=size, color=color)
        y += item_height
    return y


def card(slide, label, value, left, top, width, height,
         label_color=ACCENT, value_color=GREEN, bg=PANEL):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    shape.line.width = Pt(1)
    txt(slide, value, left, top + Inches(0.05), width, Inches(0.5),
        size=28, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    txt(slide, label, left, top + Inches(0.5), width, Inches(0.3),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)


def kv_row(slide, key, value, left, top, key_w, val_w, row_h=Inches(0.42),
           alternate=False):
    bg_color = RGBColor(0x26, 0x31, 0x47) if alternate else PANEL
    key_box = slide.shapes.add_shape(1, left, top, key_w, row_h)
    key_box.fill.solid(); key_box.fill.fore_color.rgb = bg_color
    key_box.line.color.rgb = RGBColor(0x33, 0x41, 0x55); key_box.line.width = Pt(0.5)
    val_box = slide.shapes.add_shape(1, left + key_w, top, val_w, row_h)
    val_box.fill.solid(); val_box.fill.fore_color.rgb = bg_color
    val_box.line.color.rgb = RGBColor(0x33, 0x41, 0x55); val_box.line.width = Pt(0.5)
    txt(slide, key,   left + Inches(0.1), top + Inches(0.06), key_w - Inches(0.15), row_h - Inches(0.1),
        size=12, bold=True, color=ACCENT)
    txt(slide, value, left + key_w + Inches(0.1), top + Inches(0.06), val_w - Inches(0.15), row_h - Inches(0.1),
        size=12, color=LIGHT_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

# large gradient-style centre block
hero = slide.shapes.add_shape(1, Inches(1.5), Inches(1.2), Inches(10.3), Inches(4.2))
hero.fill.solid(); hero.fill.fore_color.rgb = PANEL
hero.line.color.rgb = ACCENT; hero.line.width = Pt(2)

txt(slide, "HalluCheck v3",
    Inches(1.5), Inches(1.4), Inches(10.3), Inches(1.1),
    size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(slide, "Real-Time Hallucination Detection Middleware",
    Inches(1.5), Inches(2.55), Inches(10.3), Inches(0.6),
    size=22, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "for Large Language Models",
    Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.45),
    size=18, color=GRAY, align=PP_ALIGN.CENTER)

txt(slide, "Built with  NVIDIA NIM · ChromaDB · Redis · FastAPI · React",
    Inches(1.5), Inches(3.65), Inches(10.3), Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

# metric cards row
cw = Inches(2.3)
metrics = [
    ("KB Chunks",       "7,705+",  ACCENT),
    ("Pipeline Speed",  "~9 s",    GREEN),
    ("Test Cases",      "46 Tests",YELLOW),
    ("Project Grade",   "A",       GREEN),
]
for i, (lbl, val, col) in enumerate(metrics):
    card(slide, lbl, val,
         Inches(1.5) + i * (cw + Inches(0.15)),
         Inches(5.6), cw, Inches(0.95),
         value_color=col)

txt(slide, "AI / ML Project  —  Claim-Level Fact Verification",
    Inches(0), Inches(6.9), SLIDE_W, Inches(0.4),
    size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "01  The Problem", "Why hallucination detection matters in production")

txt(slide,
    "LLMs confidently generate false information — called hallucinations. "
    "In medicine, law, and finance, a single wrong fact causes real harm.",
    Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.55),
    size=16, color=LIGHT_TEXT)

rows = [
    ("Medical risk",   "LLM says 'Ibuprofen is safe in all trimesters' — dangerously wrong."),
    ("Legal risk",     "Wrong statute dates or fine amounts mislead professionals."),
    ("Financial risk", "Wrong FDIC limits or GDP figures cause real monetary harm."),
    ("Trust erosion",  "Users cannot tell when an LLM is fabricating — they need a safety net."),
    ("Gap in market",  "Existing tools either block everything (too strict) or do nothing (too loose)."),
]
y = Inches(2.15)
for i, (k, v) in enumerate(rows):
    kv_row(slide, k, v, Inches(0.4), y, Inches(3.2), Inches(9.0),
           row_h=Inches(0.48), alternate=(i % 2 == 1))
    y += Inches(0.48)

txt(slide,
    "There is no production-ready middleware that detects, flags, and corrects hallucinations "
    "claim-by-claim in real time — until HalluCheck.",
    Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.7),
    size=15, bold=True, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "02  Our Solution", "HalluCheck — claim-level RAG verification")

txt(slide,
    "HalluCheck sits between any LLM and the end user. "
    "It intercepts every response, extracts individual factual claims, verifies each one "
    "against an authoritative knowledge base, and rewrites wrong answers before they reach the user.",
    Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.7),
    size=15, color=LIGHT_TEXT)

# Pipeline flow
stages = [
    ("📥", "Input",    "LLM output"),
    ("🔍", "Extract",  "Claimify"),
    ("📚", "Retrieve", "KB+Web"),
    ("⚖️",  "Verify",  "LLM check"),
    ("🚦", "Decide",   "Thresholds"),
    ("✏️", "Correct",  "Rewrite"),
    ("📤", "Output",   "Safe text"),
]
box_w = Inches(1.5)
arrow_w = Inches(0.35)
total_w = len(stages)*box_w + (len(stages)-1)*arrow_w
start_x = (SLIDE_W - total_w) / 2
y_flow = Inches(2.45)
for i, (em, title, desc) in enumerate(stages):
    bx = start_x + i*(box_w + arrow_w)
    b = slide.shapes.add_shape(1, bx, y_flow, box_w, Inches(1.1))
    b.fill.solid(); b.fill.fore_color.rgb = PANEL
    b.line.color.rgb = ACCENT; b.line.width = Pt(1.2)
    txt(slide, em,    bx, y_flow + Inches(0.05), box_w, Inches(0.38), size=20, align=PP_ALIGN.CENTER, color=ACCENT)
    txt(slide, title, bx, y_flow + Inches(0.42), box_w, Inches(0.28), size=11, bold=True, color=WHITE,  align=PP_ALIGN.CENTER)
    txt(slide, desc,  bx, y_flow + Inches(0.72), box_w, Inches(0.25), size=9,  color=GRAY,  align=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        txt(slide, "→", bx + box_w, y_flow + Inches(0.3), arrow_w, Inches(0.5),
            size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

props = [
    ("Claim-level",     "Operates on individual sentences — not the whole response."),
    ("Self-correcting", "Rewrites blocked/flagged claims using verified evidence — unique feature."),
    ("Domain-aware",    "Medical needs 82% confidence; General only 40%."),
    ("Streaming SSE",   "Results stream claim-by-claim in real time — no waiting."),
    ("100% free stack", "NVIDIA NIM free tier, ChromaDB, Redis — $0/month to run."),
]
y = Inches(3.8)
for i, (k, v) in enumerate(props):
    kv_row(slide, k, v, Inches(0.4), y, Inches(3.0), Inches(9.2),
           row_h=Inches(0.42), alternate=(i % 2 == 1))
    y += Inches(0.42)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "03  Architecture", "Production-grade 6-component stack")

components = [
    ("FastAPI Backend",   "Async Python — SSE streaming, rate limiting (20 req/min), role-based auth (read + admin keys), input validation."),
    ("ChromaDB",          "Persistent vector store. HNSW cosine index. 7,705 chunks across 20+ Wikipedia articles + domain fact sheets."),
    ("Redis",             "Distributed claim cache (1h TTL). Semantic ChromaDB cache for near-duplicate queries. Prevents re-verification."),
    ("NVIDIA NIM",        "Free-tier cloud GPU inference. llama-3.2-3b (extractor, fast) + llama-3.3-70b (verifier, accurate). ~9s end-to-end."),
    ("React SPA",         "5-tab frontend (Playground, Dashboard, KB, Audit, Settings). SSE consumer for real-time claim streaming."),
    ("Docker + nginx",    "4-service Compose stack. TLS termination, HTTP→HTTPS redirect. Models baked in — zero cold-start downloads."),
]
y = Inches(1.5)
for i, (k, v) in enumerate(components):
    kv_row(slide, k, v, Inches(0.4), y, Inches(3.2), Inches(9.5),
           row_h=Inches(0.5), alternate=(i % 2 == 1))
    y += Inches(0.5)

txt(slide, "Request lifecycle:",
    Inches(0.4), Inches(5.1), Inches(6), Inches(0.35),
    size=13, bold=True, color=ACCENT)

code_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.45))
code_bg.fill.solid(); code_bg.fill.fore_color.rgb = PANEL
code_bg.line.color.rgb = RGBColor(0x33, 0x41, 0x55); code_bg.line.width = Pt(0.5)
txt(slide, "Browser → nginx:443 → middleware:8080 → [ChromaDB + Redis + NVIDIA NIM] → SSE stream back",
    Inches(0.55), Inches(5.5), Inches(12.2), Inches(0.38),
    size=12, color=GREEN)

txt(slide, "Security: Bearer token auth  |  Rate limiting  |  50K char input cap  |  URL scheme whitelist  |  HTTPS/TLS",
    Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.4),
    size=12, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PIPELINE DEEP-DIVE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "04  Pipeline Deep-Dive", "6-stage claim verification flow")

stages_detail = [
    ("Stage 1 — Claimify Extraction",
     "spaCy sentence splitting → selection (factual vs opinion) → decomposition into atomic claims → "
     "LLM produces structured JSON: text, normalized form, claim_type, stakes, category."),
    ("Stage 2 — Hybrid Retrieval (per claim)",
     "HyDE: generate a hypothetical doc to bridge phrasing gaps.  "
     "Multi-query: expand to N variants for recall.  "
     "BM25 (35%) + Vector (65%) fusion → top-10 candidates → CrossEncoder reranks → top-5 to LLM."),
    ("Stage 3 — LLM Verification",
     "llama-3.3-70b reads claim + top-5 evidence chunks.  "
     "Outputs: status (verified / contradicted / partially_supported / unverifiable), "
     "confidence 0-1, key_evidence (verbatim quote), reasoning."),
    ("Stage 4 — Domain-Aware Decision",
     "MEDICAL block<0.82  FLAG<0.88  |  LEGAL block<0.76  FLAG<0.82  |  "
     "FINANCIAL block<0.70  FLAG<0.76  |  GENERAL block<0.40  FLAG<0.60  "
     "→ action: PASS / ANNOTATE / FLAG / BLOCK"),
    ("Stage 5 — Self-Correction",
     "If any claim was blocked/flagged: a second LLM call rewrites only wrong portions "
     "using verified evidence. Arrives as a separate SSE 'corrected' event — does not delay main result."),
    ("Stage 6 — MPC Refinement (optional)",
     "Receding-horizon loop: 3 candidate phrasings per sentence → scored against KB → "
     "lowest hallucination cost selected. Off by default (adds ~3× LLM calls)."),
]
y = Inches(1.5)
for i, (title, body) in enumerate(stages_detail):
    color = ACCENT if i % 2 == 0 else YELLOW
    txt(slide, title, Inches(0.4), y, Inches(4.0), Inches(0.32),
        size=12, bold=True, color=color)
    txt(slide, body,  Inches(0.4), y + Inches(0.3), Inches(12.5), Inches(0.42),
        size=11, color=LIGHT_TEXT)
    y += Inches(0.83)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DOMAIN THRESHOLDS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "05  Domain-Aware Safety Thresholds",
            "Different stakes require different confidence requirements")

txt(slide,
    "Claims are automatically classified by the LLM into 4 categories. "
    "Each category has its own block and flag threshold — medical claims require far higher confidence than general facts.",
    Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.6),
    size=15, color=LIGHT_TEXT)

# Domain threshold cards (2×2)
domains = [
    ("MEDICAL",   "0.82", "0.88", RED,    "Drug dosages · diagnoses · treatment · clinical facts"),
    ("LEGAL",     "0.76", "0.82", YELLOW, "Regulations · court rulings · fines · dates of law"),
    ("FINANCIAL", "0.70", "0.76", YELLOW, "GDP · interest rates · market caps · tax rates"),
    ("GENERAL",   "0.40", "0.60", GREEN,  "History · geography · science · technology"),
]
card_w = Inches(3.0)
card_h = Inches(1.6)
for i, (domain, block, flag, col, examples) in enumerate(domains):
    cx = Inches(0.4) + (i % 2) * (card_w + Inches(0.35))
    cy = Inches(2.3) + (i // 2) * (card_h + Inches(0.3))
    bg_shape = slide.shapes.add_shape(1, cx, cy, card_w, card_h)
    bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = PANEL
    bg_shape.line.color.rgb = col; bg_shape.line.width = Pt(2)
    txt(slide, domain, cx, cy + Inches(0.08), card_w, Inches(0.38),
        size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(slide, f"Block < {block}  |  Flag < {flag}",
        cx, cy + Inches(0.5), card_w, Inches(0.3),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, examples,
        cx + Inches(0.1), cy + Inches(0.85), card_w - Inches(0.2), Inches(0.65),
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Example on right side
ex_x = Inches(6.8)
txt(slide, "Live Example:", ex_x, Inches(2.3), Inches(6.0), Inches(0.35),
    size=14, bold=True, color=ACCENT)

example_box = slide.shapes.add_shape(1, ex_x, Inches(2.7), Inches(6.2), Inches(1.0))
example_box.fill.solid(); example_box.fill.fore_color.rgb = PANEL
example_box.line.color.rgb = RED; example_box.line.width = Pt(1)
txt(slide,
    '"Ibuprofen is safe in all trimesters of pregnancy."',
    ex_x + Inches(0.1), Inches(2.75), Inches(6.0), Inches(0.6),
    size=13, color=WHITE)

txt(slide, "→  MEDICAL claim  →  needs confidence ≥ 0.82 to pass  →  BLOCKED (confidence 0.15)",
    ex_x, Inches(3.85), Inches(6.2), Inches(0.4),
    size=12, color=RED)

txt(slide, "→  Self-Corrector rewrites using WHO/FDA evidence from KB",
    ex_x, Inches(4.3), Inches(6.2), Inches(0.4),
    size=12, color=GREEN)

txt(slide, "→  User receives corrected, safe text — not a block error",
    ex_x, Inches(4.75), Inches(6.2), Inches(0.4),
    size=12, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KNOWLEDGE BASE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "06  Knowledge Base", "7,705+ authoritative chunks across 6 domains")

kb_rows = [
    ("Wikipedia (15+)",  "Einstein, DNA, Climate Change, Vaccination, WWII, AI, Internet, COVID-19, Diabetes, Inflation, US Constitution, Python, Stock Market, French Revolution, Quantum Mechanics, and more."),
    ("Medical facts",    "15 entries — drug safety thresholds, WHO guidelines, FDA approval dates, clinical trial results, ibuprofen pregnancy warnings."),
    ("Legal facts",      "15 entries — GDPR (Art.17/20, fines up to €20M), CCPA, Miranda rights, attorney-client privilege, double jeopardy, NDA enforceability."),
    ("Financial facts",  "15 entries — capital gains rates, Fed dual mandate, FDIC $250K limit, 401k contribution limits, Dodd-Frank, Basel III, FICO scoring."),
    ("Physics/Science",  "Speed of light (299,792,458 m/s), Planck constant, Bohr radius, laws of thermodynamics."),
    ("General tech",     "WWW invented 1989, Google founded 1998 (Mountain View), Apple co-founded 1976, iPhone launched June 2007."),
]
y = Inches(1.55)
for i, (k, v) in enumerate(kb_rows):
    kv_row(slide, k, v, Inches(0.4), y, Inches(2.8), Inches(9.9),
           row_h=Inches(0.5), alternate=(i % 2 == 1))
    y += Inches(0.5)

txt(slide, "Hybrid search config:",
    Inches(0.4), Inches(5.55), Inches(6), Inches(0.32),
    size=13, bold=True, color=ACCENT)

cfg = slide.shapes.add_shape(1, Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.38))
cfg.fill.solid(); cfg.fill.fore_color.rgb = PANEL
cfg.line.color.rgb = RGBColor(0x33, 0x41, 0x55); cfg.line.width = Pt(0.5)
txt(slide,
    "BM25_WEIGHT=0.35  |  VECTOR_WEIGHT=0.65  |  TOP_K=5  |  RERANKER=ms-marco-MiniLM-L-6-v2  |  MIN_RELEVANCE=0.20",
    Inches(0.55), Inches(5.94), Inches(12.2), Inches(0.32),
    size=11, color=GREEN)

txt(slide,
    "Web-RAG fallback: Tavily → DuckDuckGo activated when KB top score < 0.25  |  Ingestion: text / URL / Wikipedia / PDF via REST API",
    Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.4),
    size=12, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FRONTEND & UX
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "07  Frontend & User Experience", "5-tab React SPA with real-time SSE streaming")

tabs = [
    ("Playground",     "Paste any LLM output → SSE stream shows claims appearing live with VERIFIED (green) / FLAGGED (yellow) / BLOCKED (red) badges, confidence bars, verbatim evidence quotes, and self-corrected before/after diff."),
    ("Dashboard",      "Live charts: verified vs flagged vs blocked pie, confidence trend line over time, processing time bar chart, domain category breakdown, cache hit rate counter. Auto-refreshes every 10s."),
    ("Knowledge Base", "Browse all 7,705+ chunks, search by keyword, add Wikipedia topics live, upload PDFs, delete documents, see per-source chunk counts and ingest dates."),
    ("Audit Log",      "Every request logged with full claim breakdown, timestamp, model, processing time, overall confidence, block/flag/pass counts. Filter by date range. Export JSON or CSV."),
    ("Settings",       "API key management (stored in localStorage), system health check panel, cache stats + clear button, LLM provider status, KB ingest tool, backend configuration display."),
]
y = Inches(1.55)
for i, (tab, desc) in enumerate(tabs):
    kv_row(slide, tab, desc, Inches(0.4), y, Inches(2.4), Inches(10.2),
           row_h=Inches(0.52), alternate=(i % 2 == 1))
    y += Inches(0.52)

txt(slide, "Live demo text — paste this to show all 4 error types at once:",
    Inches(0.4), Inches(5.35), Inches(12.5), Inches(0.35),
    size=13, bold=True, color=ACCENT)

demo_box = slide.shapes.add_shape(1, Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.78))
demo_box.fill.solid(); demo_box.fill.fore_color.rgb = PANEL
demo_box.line.color.rgb = GREEN; demo_box.line.width = Pt(1)
txt(slide,
    '"The EU GDPR came into force on 1 June 2019. Data breaches must be reported within 24 hours. '
    'FDIC insures deposits up to $100,000. Ibuprofen is safe throughout all trimesters of pregnancy."',
    Inches(0.55), Inches(5.8), Inches(12.2), Inches(0.65),
    size=12, color=WHITE)
txt(slide, "Every single fact above is WRONG — HalluCheck blocks/flags all 4 with evidence and rewrites them.",
    Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.38),
    size=13, bold=True, color=RED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "08  Technical Stack", "Every component chosen for production readiness")

stack = [
    ("NVIDIA NIM",              "Free-tier cloud GPU. llama-3.2-3b (extract) + llama-3.3-70b (verify). ~10× faster than local Ollama."),
    ("ChromaDB",                "Persistent vector store. HNSW cosine index. Zero-cost self-hosted. GPU-accelerated embeddings."),
    ("BM25 (rank_bm25)",        "In-memory keyword index. Rebuilt from ChromaDB at startup. Adds lexical recall to semantic search."),
    ("CrossEncoder reranker",   "ms-marco-MiniLM-L-6-v2. Reranks 10 candidates → 5. Significant precision gain with <100ms overhead."),
    ("sentence-transformers",   "all-MiniLM-L6-v2. GPU-aware (CUDA auto-detected). Baked into Docker image — no cold-start download."),
    ("Redis",                   "Semantic claim cache (1h TTL). 3-level cache: semantic ChromaDB → Redis → in-memory dict."),
    ("FastAPI + uvicorn",       "Async ASGI. StreamingResponse SSE. pydantic-settings typed config. OpenAPI docs auto-generated."),
    ("React + Vite",            "SPA bundled into Docker image. No separate frontend server. recharts for live dashboard."),
    ("nginx",                   "TLS 1.2/1.3 termination. HTTP→HTTPS redirect. proxy_buffering off for SSE. Rate limit headers."),
    ("Docker Compose",          "4-service: middleware + redis + nginx. Non-root user. GPU passthrough. Offline model cache."),
    ("GitHub Actions CI",       "pytest (46 tests) + ruff lint on every push to main. Zero broken merges policy."),
]
y = Inches(1.5)
for i, (k, v) in enumerate(stack):
    kv_row(slide, k, v, Inches(0.4), y, Inches(3.0), Inches(9.6),
           row_h=Inches(0.43), alternate=(i % 2 == 1))
    y += Inches(0.43)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PERFORMANCE & METRICS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "09  Performance & Evaluation", "Benchmarks and quality metrics")

# metric cards
perf_metrics = [
    ("End-to-end (NIM)",  "~9 s",   GREEN),
    ("End-to-end (local)","80-240s",YELLOW),
    ("KB chunks",         "7,705+", ACCENT),
    ("Semantic cache hit","~100%",  GREEN),
    ("Test suite",        "46 pass",GREEN),
]
cw2 = Inches(2.35)
for i, (lbl, val, col) in enumerate(perf_metrics):
    card(slide, lbl, val,
         Inches(0.4) + i * (cw2 + Inches(0.1)),
         Inches(1.5), cw2, Inches(0.85),
         value_color=col)

txt(slide, "LLM calls per request (5 claims):",
    Inches(0.4), Inches(2.55), Inches(7), Inches(0.35),
    size=13, bold=True, color=ACCENT)

call_rows = [
    ("Claim extraction",        "1 LLM call"),
    ("Query expansion (5×)",    "5 LLM calls"),
    ("Batch verification",      "1 LLM call"),
    ("Self-correction (if any)", "1 LLM call"),
    ("Total",                   "~8 calls → ~9s on NVIDIA NIM free tier"),
]
y = Inches(2.95)
for i, (k, v) in enumerate(call_rows):
    col = GREEN if i == len(call_rows)-1 else LIGHT_TEXT
    bold_v = i == len(call_rows)-1
    kv_row(slide, k, v, Inches(0.4), y, Inches(3.6), Inches(5.8),
           row_h=Inches(0.42), alternate=(i % 2 == 1))
    y += Inches(0.42)

txt(slide, "Evaluation suite (eval_cases.yaml):",
    Inches(7.2), Inches(2.55), Inches(6.0), Inches(0.35),
    size=13, bold=True, color=ACCENT)

eval_rows = [
    ("Total test cases",     "26 handcrafted"),
    ("MEDICAL domain",       "Ibuprofen, drug dosages, diagnoses"),
    ("LEGAL domain",         "GDPR dates, Miranda rights, NDA"),
    ("FINANCIAL domain",     "FDIC limits, capital gains, Fed"),
    ("GENERAL domain",       "Wrong birthplaces, fake dates"),
    ("Evaluation command",   "python evaluate.py"),
]
y = Inches(2.95)
for i, (k, v) in enumerate(eval_rows):
    kv_row(slide, k, v, Inches(7.2), y, Inches(2.8), Inches(3.7),
           row_h=Inches(0.42), alternate=(i % 2 == 1))
    y += Inches(0.42)

txt(slide, "Cache speed-up: Second identical request served in <1ms (semantic cache hit — no LLM call).",
    Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.38),
    size=13, color=YELLOW)

txt(slide,
    "Speed options: Switch to NVIDIA NIM (~10× faster) | Reduce MAX_CLAIMS=3 | Disable MULTI_QUERY | Use smaller model",
    Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.38),
    size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — API REFERENCE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "10  API Reference", "OpenAI-compatible + custom verification endpoints")

api_rows = [
    ("POST /verify/stream",        "SSE — stream claim-by-claim results in real time (primary endpoint)."),
    ("POST /verify",               "Sync — returns complete JSON result after full pipeline runs."),
    ("POST /v1/chat/completions",  "OpenAI-compatible proxy — drop-in replacement for any OpenAI client."),
    ("POST /v1/messages",          "Anthropic-compatible proxy — hallucination detection on LLM output."),
    ("GET  /health",               "System health: LLM reachability, KB chunks, cache stats, audit totals."),
    ("GET  /audit/recent?n=50",    "Last N audit records with full claim breakdown, timestamps, confidence."),
    ("GET  /audit/stats",          "Aggregates: total verified, flagged, blocked, corrected, avg confidence."),
    ("POST /kb/ingest",            "Add text/URL to KB (admin key required). Auto-chunks and vector-indexes."),
    ("POST /kb/ingest/wikipedia",  "Fetch + chunk + index a Wikipedia article by topic name."),
    ("GET  /kb/documents",         "List all document sources with per-source chunk counts."),
    ("DELETE /kb/documents/{id}",  "Remove all chunks for a document (admin key required)."),
    ("POST /cache/clear",          "Clear Redis + in-memory + semantic cache (admin key)."),
]
y = Inches(1.5)
for i, (k, v) in enumerate(api_rows):
    kv_row(slide, k, v, Inches(0.4), y, Inches(3.6), Inches(9.0),
           row_h=Inches(0.4), alternate=(i % 2 == 1))
    y += Inches(0.4)

txt(slide, "Auth: Authorization: Bearer <key>  |  Read key: /verify /audit /health  |  Admin key: /kb/ingest /cache/clear",
    Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.38),
    size=12, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PROJECT REPORT & GRADE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "11  Honest Project Report & Grade", "Strengths, weaknesses, and final assessment")

# Left column: strengths
txt(slide, "STRENGTHS", Inches(0.4), Inches(1.5), Inches(5.8), Inches(0.35),
    size=14, bold=True, color=GREEN)
strengths = [
    "Production-grade: HTTPS, Docker, CI/CD, auth, rate limiting",
    "Self-correcting: unique feature — rewrites errors automatically",
    "Domain-aware safety with per-category thresholds",
    "3-level semantic cache — prevents redundant LLM calls",
    "Streaming SSE — live results, no blank wait screen",
    "Full observability: audit trail, dashboard, health check",
    "OpenAI + Anthropic compatible API — drop-in proxy",
    "Zero-cost stack: NVIDIA NIM free tier, ChromaDB, Redis",
    "26-case evaluation suite + 46 automated pytest tests",
    "Hybrid BM25+vector+reranker retrieval pipeline",
]
y = Inches(1.92)
for s in strengths:
    txt(slide, f"✓  {s}", Inches(0.4), y, Inches(6.0), Inches(0.33),
        size=11, color=GREEN)
    y += Inches(0.33)

# Right column: weaknesses
txt(slide, "HONEST WEAKNESSES", Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.35),
    size=14, bold=True, color=RED)
weaknesses = [
    "Rate limiting is per-process (breaks with 4 uvicorn workers) → needs Redis",
    "Sync Redis client blocks event loop — needs redis.asyncio",
    "NVIDIA NIM free tier can queue →  9–30s variance in latency",
    "English-only (spaCy en_core_web_sm) — no multilingual support",
    "No user accounts / multi-tenant isolation",
    "No automated precision/recall CI gate (eval runs manually)",
    "MPC off by default — most impressive feature not shown by default",
    "KB coverage limited to ingested topics — uningested = unverifiable",
    "No model versioning — KB drift when Chroma format changes",
    "Frontend has no pagination on audit log (can lag with 10k+ entries)",
]
y = Inches(1.92)
for w in weaknesses:
    txt(slide, f"✗  {w}", Inches(6.9), y, Inches(6.1), Inches(0.33),
        size=11, color=RED)
    y += Inches(0.33)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — GRADE BREAKDOWN
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "12  Grade Breakdown", "Category-by-category technical assessment")

grade_rows = [
    ("Architecture & Design",       "A",  "9.0/10", GREEN,
     "Clean layered separation. FastAPI + async pipeline + SSE well-designed. "
     "Minor: rate limiting and Redis client not fully async."),
    ("Feature Completeness",        "A+", "9.5/10", GREEN,
     "Self-correction, MPC, streaming, domain thresholds, Web-RAG, "
     "semantic cache, 3 LLM providers — extremely feature-rich."),
    ("Code Quality",                "A-", "8.5/10", GREEN,
     "Good Pydantic models, docstrings, error handling. "
     "Minor: some duplicate logic between verify() and verify_streaming()."),
    ("Production Readiness",        "A",  "9.0/10", GREEN,
     "HTTPS, Docker, CI, auth, rate limiting, health check, "
     "audit trail. Missing: Redis-backed rate limiting, async Redis."),
    ("Testing & Evaluation",        "B+", "8.0/10", YELLOW,
     "46 pytest tests, 26 eval cases. Missing: CI precision/recall gate, "
     "integration tests hitting real LLM, load/stress tests."),
    ("Documentation & UX",          "A",  "9.0/10", GREEN,
     "Excellent README with curl examples. 5-tab React UI. "
     "Live SSE streaming. Swagger UI. Audit log export."),
    ("Innovation",                  "A+", "9.5/10", GREEN,
     "Self-correction, domain thresholds, MPC controller, "
     "3-level semantic cache — several genuinely novel features."),
    ("Performance",                 "B+", "8.0/10", YELLOW,
     "~9s end-to-end with NVIDIA NIM is good. "
     "Semantic cache works well. Local Ollama is slow (80–240s)."),
]

# Header row
y = Inches(1.5)
for col_txt, col_w, col_x in [
    ("Category",       Inches(3.2), Inches(0.4)),
    ("Grade",          Inches(0.7), Inches(3.6)),
    ("Score",          Inches(0.9), Inches(4.3)),
    ("Notes",          Inches(8.0), Inches(5.2)),
]:
    hdr = slide.shapes.add_shape(1, col_x, y, col_w, Inches(0.38))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    hdr.line.color.rgb = ACCENT; hdr.line.width = Pt(1)
    txt(slide, col_txt, col_x + Inches(0.05), y + Inches(0.05), col_w, Inches(0.3),
        size=11, bold=True, color=ACCENT)

y += Inches(0.38)
for i, (cat, grade, score, col, notes) in enumerate(grade_rows):
    bg = RGBColor(0x26, 0x31, 0x47) if i % 2 == 1 else PANEL
    row_h = Inches(0.5)
    for shape_x, shape_w in [
        (Inches(0.4), Inches(3.2)),
        (Inches(3.6), Inches(0.7)),
        (Inches(4.3), Inches(0.9)),
        (Inches(5.2), Inches(8.0)),
    ]:
        b = slide.shapes.add_shape(1, shape_x, y, shape_w, row_h)
        b.fill.solid(); b.fill.fore_color.rgb = bg
        b.line.color.rgb = RGBColor(0x33, 0x41, 0x55); b.line.width = Pt(0.4)
    txt(slide, cat,   Inches(0.5),  y + Inches(0.1), Inches(3.0), row_h, size=11, color=LIGHT_TEXT)
    txt(slide, grade, Inches(3.62), y + Inches(0.05), Inches(0.65), row_h, size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(slide, score, Inches(4.35), y + Inches(0.1), Inches(0.8), row_h, size=11, color=col, align=PP_ALIGN.CENTER)
    txt(slide, notes, Inches(5.3),  y + Inches(0.08), Inches(7.8), row_h, size=10, color=GRAY)
    y += row_h

# Overall grade
ov_box = slide.shapes.add_shape(1, Inches(0.4), y + Inches(0.1), Inches(12.5), Inches(0.55))
ov_box.fill.solid(); ov_box.fill.fore_color.rgb = RGBColor(0x05, 0x20, 0x10)
ov_box.line.color.rgb = GREEN; ov_box.line.width = Pt(2)
txt(slide, "OVERALL GRADE:  A   (8.9 / 10.0)   —   Production-ready AI safety middleware with genuine innovation",
    Inches(0.55), y + Inches(0.15), Inches(12.2), Inches(0.38),
    size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — WHAT TO IMPROVE NEXT
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
heading_bar(slide, "13  What to Build Next", "Roadmap to make the project truly world-class")

roadmap = [
    ("HIGH IMPACT — Tier 1",  ACCENT, [
        "Redis-backed rate limiting (ZADD/ZREMRANGEBYSCORE) — fixes multi-worker gap",
        "Async Redis client (redis.asyncio) — unblocks event loop on cache ops",
        "Streaming claim extraction — pipe claims to verifier as they are extracted",
        "Async KB ingest — offload upsert() to thread so large docs don't stall the loop",
    ]),
    ("QUALITY — Tier 2",  YELLOW, [
        "CI precision/recall gate — fail the build if accuracy drops below threshold",
        "Confidence calibration (Platt scaling) — raw LLM confidence is overconfident",
        "Per-claim batch isolation — stop one critical claim from forcing full ensemble",
        "Multilingual support — add spaCy xx_ent_wiki_sm for non-English KB ingestion",
    ]),
    ("INFRASTRUCTURE — Tier 3",  GREEN, [
        "Prometheus /metrics endpoint + Grafana dashboard (latency p95, cache hit rate)",
        "Model versioning — version-stamp embeddings so KB can survive Chroma upgrades",
        "Multi-tenant isolation — per-user KB namespaces + per-user audit logs",
        "Automated evaluation in CI — run eval_cases.yaml on every pull request",
    ]),
]

y = Inches(1.5)
for tier_name, col, items in roadmap:
    txt(slide, tier_name, Inches(0.4), y, Inches(12.5), Inches(0.32),
        size=13, bold=True, color=col)
    y += Inches(0.32)
    for item in items:
        txt(slide, f"  →  {item}", Inches(0.5), y, Inches(12.2), Inches(0.33),
            size=11, color=LIGHT_TEXT)
        y += Inches(0.33)
    y += Inches(0.12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

hero2 = slide.shapes.add_shape(1, Inches(1.5), Inches(1.0), Inches(10.3), Inches(4.5))
hero2.fill.solid(); hero2.fill.fore_color.rgb = PANEL
hero2.line.color.rgb = ACCENT; hero2.line.width = Pt(2)

txt(slide, "HalluCheck v3",
    Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.0),
    size=52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(slide, "Keeping AI Honest — One Claim at a Time",
    Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.55),
    size=22, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "Quick-start:  python run_proxy.py  →  http://localhost:8080",
    Inches(1.5), Inches(3.05), Inches(10.3), Inches(0.45),
    size=16, color=GREEN, align=PP_ALIGN.CENTER)

# final metric cards
final_metrics = [
    ("Extractor",  "llama-3.2-3b",  ACCENT),
    ("Verifier",   "llama-3.3-70b", ACCENT),
    ("Vector DB",  "ChromaDB",      YELLOW),
    ("Cache",      "Redis",         YELLOW),
    ("Frontend",   "React+Vite",    GREEN),
    ("Grade",      "A",             GREEN),
]
cw3 = Inches(1.9)
for i, (lbl, val, col) in enumerate(final_metrics):
    card(slide, lbl, val,
         Inches(0.9) + i * (cw3 + Inches(0.1)),
         Inches(5.55), cw3, Inches(0.9),
         value_color=col)

txt(slide, "Thank you for reviewing HalluCheck v3",
    Inches(0), Inches(6.7), SLIDE_W, Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
OUTPUT = "HalluCheck_v3_Presentation.pptx"
prs.save(OUTPUT)
print(f"\nPresentation saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
